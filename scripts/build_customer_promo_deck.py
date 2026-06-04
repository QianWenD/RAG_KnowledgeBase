# -*- coding: utf-8 -*-
"""Build customer-facing PPTX and HTML decks for the promotion kit.

The script reads UTF-8 Markdown from disk and writes PowerPoint text through
python-pptx Unicode APIs. This avoids the command-line encoding issue that can
turn Chinese characters into question marks on Windows consoles.
"""

from __future__ import annotations

from dataclasses import dataclass
import argparse
import html
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
KIT_DIR = ROOT / "docs" / "customer-promo-kit"
SOURCE = KIT_DIR / "customer-demo-deck.md"
PPTX_OUTPUT = KIT_DIR / "rag-knowledge-center-customer-deck.pptx"
HTML_OUTPUT = KIT_DIR / "rag-knowledge-center-customer-deck.html"

FONT_CN = "Microsoft YaHei"
FONT_TITLE = "DengXian"

COLORS = {
    "ink": RGBColor(15, 36, 53),
    "deep": RGBColor(5, 28, 42),
    "teal": RGBColor(14, 126, 118),
    "mint": RGBColor(210, 241, 232),
    "gold": RGBColor(176, 123, 45),
    "paper": RGBColor(248, 245, 236),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(218, 225, 219),
    "muted": RGBColor(105, 122, 132),
    "danger": RGBColor(184, 72, 58),
}


@dataclass
class Slide:
    title: str
    subtitle: str
    bullets: list[str]
    note: str = ""


def clean_inline(text: str) -> str:
    text = text.strip()
    text = text.replace("`", "")
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    return text


def parse_markdown(source: Path) -> list[Slide]:
    markdown = source.read_text(encoding="utf-8")
    sections = [section.strip() for section in re.split(r"(?m)^---\s*$", markdown) if section.strip()]
    slides: list[Slide] = []

    for section in sections:
        lines = section.splitlines()
        heading = ""
        for line in lines:
            if line.startswith("## "):
                heading = clean_inline(line[3:])
                break
            if line.startswith("# "):
                heading = clean_inline(line[2:])
                break
        if not heading:
            continue

        title = re.sub(r"^\d+\.\s*", "", heading).strip()
        subtitle = ""
        bullets: list[str] = []
        note_parts: list[str] = []
        in_code = False
        capture_note = False

        for line in lines:
            stripped = line.strip()
            if stripped.startswith("```"):
                in_code = not in_code
                continue
            if in_code:
                continue
            if not stripped:
                continue
            if stripped.startswith("#"):
                continue
            if stripped.endswith("：") and not stripped.startswith("- "):
                capture_note = stripped in {"讲解重点：", "演示建议：", "讲解提醒："}
                continue
            if stripped.startswith(">"):
                subtitle = clean_inline(stripped.lstrip("> "))
                continue
            if stripped.startswith("- "):
                item = clean_inline(stripped[2:])
                if capture_note and len(note_parts) < 2:
                    note_parts.append(item)
                elif len(bullets) < 5:
                    bullets.append(item)
                continue
            if not subtitle and len(stripped) <= 42:
                subtitle = clean_inline(stripped)

        if title == "封面":
            title = "RAG知识中心平台"
            subtitle = "让企业资料变成可权限管控、可追溯、可运营的智能问答能力"
        if not bullets:
            bullets = note_parts[:4]
        if title and bullets:
            slides.append(Slide(title=title, subtitle=subtitle, bullets=bullets, note="；".join(note_parts)))

    return slides


def set_run_font(run, *, size: float, color: str, bold: bool = False, title_font: bool = False) -> None:
    font_name = FONT_TITLE if title_font else FONT_CN
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    r_pr = run._r.get_or_add_rPr()
    r_pr.set("lang", "zh-CN")
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            r_pr.append(el)
        el.set("typeface", font_name)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float, color: str, bold: bool = False, align=PP_ALIGN.LEFT, title_font: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, color=color, bold=bold, title_font=title_font)
    return box


def add_background(slide, page: int) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["paper"]

    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1.04), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = COLORS["deep"]
    side.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.04), Inches(0), Inches(12.3), Inches(0.16))
    top.fill.solid()
    top.fill.fore_color.rgb = COLORS["teal"]
    top.line.fill.background()

    gold = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.04), Inches(7.24), Inches(12.3), Inches(0.06))
    gold.fill.solid()
    gold.fill.fore_color.rgb = COLORS["gold"]
    gold.line.fill.background()

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.24), Inches(0.35), Inches(0.56), Inches(0.56))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS["teal"]
    circle.line.color.rgb = COLORS["mint"]
    add_text(slide, f"{page:02d}", 0.19, 6.84, 0.65, 0.32, size=11, color="white", bold=True, align=PP_ALIGN.CENTER)


def add_title(slide, item: Slide, page: int) -> None:
    add_background(slide, page)
    add_text(slide, item.title, 1.42, 0.58, 8.7, 0.64, size=28, color="ink", bold=True, title_font=True)
    if item.subtitle:
        add_text(slide, item.subtitle, 1.45, 1.21, 9.4, 0.36, size=11.8, color="muted")
    add_text(slide, "RAG知识中心平台", 10.2, 0.66, 2.35, 0.28, size=10.5, color="teal", bold=True, align=PP_ALIGN.RIGHT)


def add_bullet_panel(slide, bullets: list[str], x: float, y: float, w: float, h: float) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["white"]
    panel.line.color.rgb = COLORS["line"]
    panel.line.width = Pt(0.8)

    top = y + 0.36
    for index, bullet in enumerate(bullets):
        cy = top + index * 0.78
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.35), Inches(cy + 0.05), Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLORS["teal"] if index % 2 == 0 else COLORS["gold"]
        marker.line.fill.background()
        add_text(slide, bullet, x + 0.68, cy - 0.04, w - 1.05, 0.5, size=15.2, color="ink")


def add_side_card(slide, title: str, text: str) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.62), Inches(1.78), Inches(3.62), Inches(4.6))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["mint"]
    card.line.color.rgb = COLORS["line"]
    add_text(slide, title, 8.95, 2.15, 2.92, 0.36, size=14, color="teal", bold=True)
    add_text(slide, text, 8.95, 2.75, 2.82, 2.75, size=14, color="ink")


def add_flow_slide(slide, item: Slide, page: int) -> None:
    add_title(slide, item, page)
    steps = ["资料准备", "上传入库", "解析切片", "向量索引", "权限绑定", "用户提问", "检索授权知识", "生成回答", "引用追溯"]
    x0, y0 = 1.42, 2.02
    for index, step in enumerate(steps):
        row = index // 3
        col = index % 3
        x = x0 + col * 3.42
        y = y0 + row * 1.34
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.55), Inches(0.68))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS["white"]
        box.line.color.rgb = COLORS["teal"] if index < 5 else COLORS["gold"]
        add_text(slide, f"{index + 1}. {step}", x + 0.18, y + 0.17, 2.18, 0.3, size=13.2, color="ink", bold=True, align=PP_ALIGN.CENTER)
        if col < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.62), Inches(y + 0.22), Inches(0.42), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["teal"]
            arrow.line.fill.background()
    add_side_card(slide, "客户理解重点", "资料进入系统后，先被解析和索引，再按权限进入问答范围。用户看到的答案来自被授权的知识，而不是无边界生成。")


def add_cover(slide, item: Slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["deep"]

    band = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.4), Inches(0), Inches(5.95), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["paper"]
    band.line.fill.background()

    halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(0.72), Inches(3.5), Inches(3.5))
    halo.fill.solid()
    halo.fill.fore_color.rgb = COLORS["mint"]
    halo.line.fill.background()

    add_text(slide, "ENTERPRISE KNOWLEDGE OPS", 1.26, 1.62, 4.0, 0.28, size=10, color="gold", bold=True)
    add_text(slide, item.title, 1.24, 2.06, 5.8, 1.36, size=38, color="white", bold=True, title_font=True)
    add_text(slide, item.subtitle, 1.28, 3.64, 5.78, 0.56, size=15, color="white")
    add_bullet_panel(slide, item.bullets[:3], 8.12, 2.08, 3.85, 2.72)
    add_text(slide, "客户讲解版", 8.16, 5.56, 2.0, 0.28, size=12, color="teal", bold=True)
    add_text(slide, "适用于产品介绍、试点启动、上线培训前说明", 8.16, 5.95, 3.78, 0.38, size=11.2, color="muted")


def add_regular_slide(slide, item: Slide, page: int) -> None:
    add_title(slide, item, page)
    add_bullet_panel(slide, item.bullets, 1.42, 1.86, 6.78, 4.72)
    side_text = item.note or "先讲业务价值，再讲系统能力。客户最关心的是能不能上线、能不能控权限、出了问题能不能追溯。"
    add_side_card(slide, "讲解提示", side_text)


def build_pptx(slides: list[Slide], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for page, item in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if page == 1:
            add_cover(slide, item)
        elif "整体业务流程" in item.title:
            add_flow_slide(slide, item, page)
        else:
            add_regular_slide(slide, item, page)

    prs.save(output)


def render_html(slides: list[Slide], output: Path) -> None:
    slide_sections = []
    for index, slide in enumerate(slides, start=1):
        bullets = "\n".join(f"<li>{html.escape(item)}</li>" for item in slide.bullets)
        flow = ""
        if "整体业务流程" in slide.title:
            steps = ["资料准备", "上传入库", "解析切片", "向量索引", "权限绑定", "用户提问", "检索授权知识", "生成回答", "引用追溯"]
            flow = "<div class=\"flow\">" + "".join(f"<span>{html.escape(step)}</span>" for step in steps) + "</div>"
        slide_sections.append(
            f"""
            <section class="slide {'cover' if index == 1 else ''}">
              <div class="index">{index:02d}</div>
              <div class="brand">RAG知识中心平台</div>
              <h1>{html.escape(slide.title)}</h1>
              <p class="subtitle">{html.escape(slide.subtitle)}</p>
              {flow}
              <ul>{bullets}</ul>
              <aside>{html.escape(slide.note or '客户要记住：可部署、可授权、可追溯、可运营。')}</aside>
            </section>
            """
        )

    page = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>RAG知识中心平台客户讲解版</title>
  <style>
    :root {{
      --deep: #051c2a;
      --ink: #0f2435;
      --teal: #0e7e76;
      --mint: #d2f1e8;
      --gold: #b07b2d;
      --paper: #f8f5ec;
      --line: #dae1db;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: #d8d3c5;
      color: var(--ink);
      font-family: "Microsoft YaHei", "DengXian", "Noto Sans CJK SC", sans-serif;
    }}
    .deck {{
      display: grid;
      gap: 28px;
      padding: 28px;
      justify-items: center;
    }}
    .slide {{
      position: relative;
      width: min(1280px, 96vw);
      aspect-ratio: 16 / 9;
      overflow: hidden;
      border-radius: 24px;
      padding: 72px 96px 64px 152px;
      background:
        radial-gradient(circle at 78% 18%, rgba(210, 241, 232, .95), transparent 23%),
        linear-gradient(135deg, #fffdf7 0%, var(--paper) 60%, #ece5d6 100%);
      box-shadow: 0 28px 80px rgba(5, 28, 42, .18);
    }}
    .slide::before {{
      content: "";
      position: absolute;
      inset: 0 auto 0 0;
      width: 96px;
      background: var(--deep);
    }}
    .slide::after {{
      content: "";
      position: absolute;
      left: 96px;
      right: 0;
      bottom: 28px;
      height: 6px;
      background: linear-gradient(90deg, var(--teal), var(--gold));
      border-radius: 999px;
    }}
    .cover {{
      padding-left: 112px;
      background:
        linear-gradient(112deg, var(--deep) 0%, var(--deep) 52%, transparent 52.2%),
        radial-gradient(circle at 78% 35%, var(--mint), transparent 22%),
        linear-gradient(135deg, #fffdf7, #e9dfca);
    }}
    .cover::before {{ display: none; }}
    .brand {{
      position: absolute;
      right: 92px;
      top: 54px;
      color: var(--teal);
      font-weight: 800;
      letter-spacing: .12em;
      font-size: 15px;
    }}
    .index {{
      position: absolute;
      left: 25px;
      bottom: 44px;
      z-index: 1;
      color: white;
      font-weight: 800;
      letter-spacing: .1em;
    }}
    h1 {{
      margin: 0;
      max-width: 760px;
      font-size: clamp(36px, 5vw, 62px);
      line-height: 1.08;
      letter-spacing: -.04em;
      font-family: "DengXian", "Microsoft YaHei", sans-serif;
    }}
    .cover h1, .cover .subtitle {{
      color: white;
      max-width: 560px;
    }}
    .subtitle {{
      margin: 18px 0 30px;
      max-width: 780px;
      color: #60727e;
      font-size: clamp(18px, 2vw, 24px);
      line-height: 1.55;
    }}
    ul {{
      width: min(720px, 64%);
      margin: 0;
      padding: 28px 36px 28px 54px;
      border: 1px solid var(--line);
      border-radius: 24px;
      background: rgba(255, 255, 255, .86);
      font-size: clamp(18px, 2vw, 25px);
      line-height: 1.72;
    }}
    li::marker {{ color: var(--teal); }}
    aside {{
      position: absolute;
      right: 84px;
      top: 258px;
      width: 300px;
      min-height: 210px;
      padding: 26px;
      border-radius: 24px;
      background: var(--mint);
      color: var(--ink);
      font-size: 20px;
      line-height: 1.7;
      box-shadow: inset 0 0 0 1px rgba(14, 126, 118, .12);
    }}
    .cover aside {{ display: none; }}
    .cover ul {{
      position: absolute;
      right: 88px;
      top: 220px;
      width: 380px;
      color: var(--ink);
    }}
    .flow {{
      display: grid;
      grid-template-columns: repeat(3, 1fr);
      gap: 18px;
      width: 740px;
      margin-top: 24px;
    }}
    .flow span {{
      display: grid;
      place-items: center;
      min-height: 74px;
      border-radius: 20px;
      background: white;
      border: 2px solid var(--teal);
      font-size: 22px;
      font-weight: 800;
    }}
    .flow + ul {{ display: none; }}
    @media print {{
      body {{ background: white; }}
      .deck {{ padding: 0; gap: 0; }}
      .slide {{ width: 100vw; border-radius: 0; box-shadow: none; page-break-after: always; }}
    }}
  </style>
</head>
<body>
  <main class="deck">
    {''.join(slide_sections)}
  </main>
</body>
</html>
"""
    output.write_text(page, encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build customer promotion deck assets.")
    parser.add_argument("--source", type=Path, default=SOURCE)
    parser.add_argument("--pptx", type=Path, default=PPTX_OUTPUT)
    parser.add_argument("--html", type=Path, default=HTML_OUTPUT)
    args = parser.parse_args()

    source = args.source if args.source.is_absolute() else ROOT / args.source
    pptx = args.pptx if args.pptx.is_absolute() else ROOT / args.pptx
    html_output = args.html if args.html.is_absolute() else ROOT / args.html

    slides = parse_markdown(source)
    if len(slides) < 10:
        raise SystemExit(f"Expected at least 10 slides, got {len(slides)}")

    build_pptx(slides, pptx)
    render_html(slides, html_output)
    print(f"PPTX: {pptx}")
    print(f"HTML: {html_output}")


if __name__ == "__main__":
    main()
