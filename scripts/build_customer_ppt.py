# -*- coding: utf-8 -*-
"""Build the customer-facing PowerPoint deck from the Markdown script.

The previous one-off generation command embedded Chinese text directly in the
PowerShell command body, which can pass through a non-UTF-8 code page and turn
Chinese characters into question marks. This script reads the UTF-8 Markdown
file from disk instead, keeping the generated PPTX reproducible.
"""

from __future__ import annotations

from dataclasses import dataclass
import argparse
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "docs" / "客户讲解PPT文稿.md"
OUTPUT = ROOT / "docs" / "客户讲解PPT草稿.pptx"

FONT_NAME = "Microsoft YaHei"

COLORS = {
    "navy": RGBColor(20, 43, 67),
    "blue": RGBColor(24, 119, 242),
    "cyan": RGBColor(32, 180, 190),
    "green": RGBColor(64, 150, 105),
    "orange": RGBColor(236, 139, 64),
    "bg": RGBColor(245, 248, 250),
    "white": RGBColor(255, 255, 255),
    "text": RGBColor(38, 52, 63),
    "muted": RGBColor(101, 118, 132),
    "line": RGBColor(218, 226, 232),
}


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: list[str]
    note: str = ""
    demo: list[str] | None = None


def strip_markdown_inline(text: str) -> str:
    text = text.strip()
    text = text.replace("`", "")
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    return text


def parse_slide_sections(markdown: str) -> list[SlideSpec]:
    sections = re.split(r"(?m)^## ", markdown)
    slides: list[SlideSpec] = []

    for raw in sections[1:]:
        lines = raw.splitlines()
        if not lines:
            continue

        heading = lines[0].strip()
        if not heading.startswith("第 "):
            continue

        section_title = heading.split("：", 1)[-1].strip()
        title = section_title
        bullets: list[str] = []
        note = ""
        demo: list[str] = []
        mode = ""

        for line in lines[1:]:
            stripped = line.strip()
            if stripped.startswith("标题："):
                title = strip_markdown_inline(stripped.split("：", 1)[1])
                continue
            if stripped == "要点：":
                mode = "bullets"
                continue
            if stripped == "讲解备注：":
                mode = "note"
                continue
            if stripped == "演示动作：":
                mode = "demo"
                continue
            if stripped.endswith("建议：") or stripped.endswith("流程：") or stripped.endswith("时长："):
                mode = ""
                continue
            if not stripped:
                continue
            if stripped.startswith("```"):
                mode = ""
                continue
            if stripped.startswith("- "):
                item = strip_markdown_inline(stripped[2:])
                if mode == "bullets":
                    bullets.append(item)
                elif mode == "demo":
                    demo.append(item)
                continue
            if mode == "note" and not note and not stripped.startswith("|"):
                note = strip_markdown_inline(stripped)

        if title and bullets:
            slides.append(SlideSpec(title=title, subtitle="", bullets=bullets[:5], note=note, demo=demo[:5] or None))

    return slides[:12]


def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, *, size: float, color: str, bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333333), Inches(0.22))
    top.fill.solid()
    top.fill.fore_color.rgb = COLORS["blue"]
    top.line.fill.background()

    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.22), Inches(13.333333), Inches(0.28))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = COLORS["navy"]
    bottom.line.fill.background()


def add_header(slide, title: str, page: int):
    add_textbox(slide, title, 0.65, 0.56, 9.4, 0.48, size=25, color="navy", bold=True)
    add_textbox(slide, f"{page:02d}", 12.05, 0.56, 0.75, 0.35, size=14, color="muted", bold=True, align=PP_ALIGN.RIGHT)


def add_bullets(slide, bullets: list[str], x=0.78, y=1.55, w=6.05, h=4.8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = FONT_NAME
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS["text"]


def add_card(slide, title: str, items: list[str], x=7.25, y=1.55, w=5.15, h=4.8, color="cyan"):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["line"]
    card.line.width = Pt(1)

    add_textbox(slide, title, x + 0.35, y + 0.32, w - 0.7, 0.34, size=15, color="navy", bold=True)
    top = y + 0.9
    for index, item in enumerate(items[:5]):
        cy = top + index * 0.72
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.36), Inches(cy + 0.05), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLORS[color]
        dot.line.fill.background()
        add_textbox(slide, item, x + 0.65, cy - 0.02, w - 1.0, 0.45, size=12.5, color="text")


def add_cover(slide, spec: SlideSpec):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["navy"]

    accent = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(8.7), Inches(-1.2), Inches(5.5), Inches(5.5))
    accent.line.color.rgb = COLORS["cyan"]
    accent.line.width = Pt(5)
    accent.fill.background()

    add_textbox(slide, spec.title, 0.85, 1.65, 9.2, 0.8, size=34, color="white", bold=True)
    add_textbox(slide, "客户讲解版 | 知识上传、权限控制、智能问答、审计运维", 0.9, 2.55, 8.8, 0.42, size=16, color="white")
    add_card(slide, "核心定位", spec.bullets[:3], x=0.95, y=3.45, w=5.6, h=2.35, color="cyan")
    add_textbox(slide, "客户讲解草稿", 9.25, 5.85, 3.0, 0.3, size=14, color="white", bold=True, align=PP_ALIGN.RIGHT)
    add_textbox(slide, "2026-05-25", 9.25, 6.25, 3.0, 0.3, size=12, color="white", align=PP_ALIGN.RIGHT)


def build_deck(slides: list[SlideSpec], output: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for page, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if page == 1:
            add_cover(slide, spec)
            continue

        add_background(slide)
        add_header(slide, spec.title, page)
        add_bullets(slide, spec.bullets)

        if spec.demo:
            add_card(slide, "演示动作", spec.demo, color="orange")
        elif spec.note:
            add_card(slide, "讲解提示", [spec.note], color="green")
        else:
            add_card(slide, "客户要记住", spec.bullets[:3], color="cyan")

    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the customer-facing PowerPoint deck.")
    parser.add_argument("--source", type=Path, default=SOURCE, help="Markdown source file")
    parser.add_argument("--output", type=Path, default=OUTPUT, help="PPTX output file")
    args = parser.parse_args()

    source = args.source if args.source.is_absolute() else ROOT / args.source
    output = args.output if args.output.is_absolute() else ROOT / args.output

    markdown = source.read_text(encoding="utf-8")
    slides = parse_slide_sections(markdown)
    if len(slides) < 10:
        raise SystemExit(f"Expected at least 10 slide sections, got {len(slides)}")
    build_deck(slides, output)
    print(output)


if __name__ == "__main__":
    main()
